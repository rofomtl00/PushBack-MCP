"""
parser.py — Multi-Format Document Parser
==========================================
Extracts text, tables, and metadata from business documents.
Supports: PDF, DOCX, XLSX, PPTX, CSV, TXT, images.
"""

import os
import csv
import json


MAX_PARSE_SIZE_MB = 30  # Skip parsing files larger than this (binary content bloats memory)

def parse_file(filepath: str) -> dict:
    """Parse a file and return structured content."""
    ext = os.path.splitext(filepath)[1].lower()
    name = os.path.basename(filepath)
    size_kb = round(os.path.getsize(filepath) / 1024, 1)

    result = {
        "filename": name,
        "type": ext,
        "size_kb": size_kb,
        "text": "",
        "tables": [],
        "metadata": {},
        "error": None,
    }

    try:
        if ext == ".pdf":
            result.update(_parse_pdf(filepath))
        elif ext == ".docx":
            result.update(_parse_docx(filepath))
        elif ext in (".xlsx", ".xls"):
            result.update(_parse_xlsx(filepath))
        elif ext == ".pptx":
            result.update(_parse_pptx(filepath))
        elif ext == ".csv":
            result.update(_parse_csv(filepath))
        elif ext in (".txt", ".md", ".py", ".js", ".ts", ".tsx", ".jsx", ".go",
                      ".rs", ".java", ".cpp", ".c", ".h", ".rb", ".php", ".swift",
                      ".kt", ".html", ".css", ".scss", ".sql", ".sh", ".json",
                      ".yaml", ".yml", ".toml", ".fdx", ".fountain"):
            with open(filepath, "r", errors="ignore") as f:
                result["text"] = f.read(100000)
        elif ext in (".png", ".jpg", ".jpeg", ".gif", ".webp", ".tiff", ".bmp", ".svg"):
            result["metadata"]["type"] = "image"
            result["metadata"]["format"] = ext
            # Try to get image dimensions
            try:
                from PIL import Image
                img = Image.open(filepath)
                result["metadata"]["dimensions"] = f"{img.width}x{img.height}"
            except Exception:
                pass
            # Try OCR text extraction
            image_text = ""
            try:
                from PIL import Image
                import pytesseract
                img = Image.open(filepath)
                image_text = pytesseract.image_to_string(img).strip()
            except ImportError:
                image_text = "[Image analysis requires: pip install pytesseract Pillow]"
            except Exception as e:
                image_text = f"[Could not read image: {e}]"
            if image_text and image_text != "[Image analysis requires: pip install pytesseract Pillow]":
                result["text"] = f"[Image: {name}]\nExtracted text:\n{image_text[:5000]}"
            else:
                result["text"] = f"[Image: {name}] {image_text}"
        # Video files — can't read content but log metadata
        elif ext in (".mp4", ".mov", ".avi", ".mkv", ".webm", ".flv", ".wmv", ".m4v",
                      ".prproj", ".drp", ".fcpxml", ".aep"):
            size_mb = round(size_kb / 1024, 1)
            result["text"] = f"[Video/Project file: {name} ({size_mb} MB)]"
            result["metadata"]["type"] = "video" if ext not in (".prproj", ".drp", ".fcpxml", ".aep") else "video-project"
            result["metadata"]["size_mb"] = size_mb
        # Audio files
        elif ext in (".mp3", ".wav", ".flac", ".aac", ".ogg", ".m4a", ".aif", ".aiff",
                      ".als", ".flp", ".logic", ".ptx", ".rpp"):
            size_mb = round(size_kb / 1024, 1)
            result["text"] = f"[Audio/Music file: {name} ({size_mb} MB)]"
            result["metadata"]["type"] = "audio" if ext not in (".als", ".flp", ".logic", ".ptx", ".rpp") else "music-project"
            result["metadata"]["size_mb"] = size_mb
        # 3D/Graphics files
        elif ext in (".blend", ".fbx", ".obj", ".glb", ".gltf", ".stl", ".usd",
                      ".c4d", ".max", ".ma", ".psd", ".ai", ".xcf", ".sketch", ".fig",
                      ".indd", ".afdesign", ".afphoto"):
            size_mb = round(size_kb / 1024, 1)
            project_types = {".blend":"Blender", ".c4d":"Cinema 4D", ".max":"3ds Max", ".ma":"Maya",
                           ".psd":"Photoshop", ".ai":"Illustrator", ".xcf":"GIMP", ".sketch":"Sketch",
                           ".fig":"Figma", ".indd":"InDesign", ".afdesign":"Affinity Designer",
                           ".afphoto":"Affinity Photo"}
            app_name = project_types.get(ext, "3D/Graphics")
            result["text"] = f"[{app_name} project: {name} ({size_mb} MB)]"
            result["metadata"]["type"] = "graphics-project"
            result["metadata"]["application"] = app_name
            result["metadata"]["size_mb"] = size_mb
        # Medical/scientific files
        elif ext in (".dcm", ".nii", ".nii.gz", ".dicom"):
            result["text"] = f"[Medical imaging: {name} (DICOM)]"
            result["metadata"]["type"] = "medical-imaging"
        # CAD/Engineering
        elif ext in (".dwg", ".dxf", ".step", ".stp", ".iges", ".igs"):
            result["text"] = f"[CAD/Engineering: {name}]"
            result["metadata"]["type"] = "cad"
        else:
            # Try to read as text, fall back to metadata only
            try:
                with open(filepath, "r", errors="ignore") as f:
                    text = f.read(100000)
                if text and len(text.strip()) > 10 and not any(c in text[:200] for c in ['\x00', '\xff', '\xfe']):
                    result["text"] = text
                else:
                    size_mb = round(size_kb / 1024, 1)
                    result["text"] = f"[Binary file: {name} ({size_mb} MB)]"
                    result["metadata"]["type"] = "binary"
            except Exception:
                result["text"] = f"[File: {name} ({size_kb} KB)]"
    except MemoryError:
        result["error"] = "File too large to process in memory"
        result["text"] = f"[{name}: too large to parse — {size_kb} KB]"
    except Exception as e:
        result["error"] = str(e)

    # Safety: cap extracted text to prevent memory bloat from single files
    if len(result.get("text", "")) > 500_000:
        result["text"] = result["text"][:500_000] + "\n[Content truncated at 500K chars for this file]"

    return result


def _parse_pdf(filepath):
    try:
        from PyPDF2 import PdfReader
    except ImportError:
        return {"text": "", "error": "PyPDF2 not installed"}

    reader = PdfReader(filepath)
    pages = []
    for i, page in enumerate(reader.pages[:50]):  # Max 50 pages
        text = page.extract_text() or ""
        if text.strip():
            pages.append(text)

    metadata = {}
    if reader.metadata:
        for key in ("title", "author", "subject", "creator"):
            val = getattr(reader.metadata, key, None)
            if val:
                metadata[key] = str(val)

    return {
        "text": "\n\n--- Page Break ---\n\n".join(pages),
        "metadata": {**metadata, "pages": len(reader.pages)},
    }


def _parse_docx(filepath):
    try:
        from docx import Document
    except ImportError:
        return {"text": "", "error": "python-docx not installed"}

    doc = Document(filepath)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    text = "\n".join(paragraphs)

    # Extract tables
    tables = []
    for table in doc.tables[:10]:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(cells)
        if rows:
            tables.append(rows)

    return {"text": text, "tables": tables}


def _parse_xlsx(filepath):
    try:
        from openpyxl import load_workbook
    except ImportError:
        return {"text": "", "error": "openpyxl not installed"}

    wb = load_workbook(filepath, data_only=True)
    sheets = []
    tables = []

    for ws in wb.worksheets[:10]:
        rows = []
        for row in ws.iter_rows(max_row=200, values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(c.strip() for c in cells):
                rows.append(cells)
        if rows:
            sheets.append({"sheet": ws.title, "rows": len(rows)})
            tables.append(rows)

    # Formula extraction: two-pass approach
    formulas_found = []
    try:
        wb_formulas = load_workbook(filepath, data_only=False)
        wb_values = load_workbook(filepath, data_only=True)
        all_sheet_names = set(wb_formulas.sheetnames)

        for sheet_name in wb_formulas.sheetnames[:10]:
            ws_f = wb_formulas[sheet_name]
            ws_v = wb_values[sheet_name]
            for row in ws_f.iter_rows(max_row=500):
                for cell in row:
                    if cell.value and str(cell.value).startswith('='):
                        if len(formulas_found) >= 100:
                            break
                        formula_text = str(cell.value)
                        val_cell = ws_v[cell.coordinate]
                        computed = str(val_cell.value) if val_cell.value is not None else "N/A"

                        # Flag potential issues
                        issues = []
                        import re
                        # Check for references to non-existent sheets
                        sheet_refs = re.findall(r"'?([A-Za-z0-9_ ]+)'?!", formula_text)
                        for ref in sheet_refs:
                            if ref not in all_sheet_names:
                                issues.append(f"references non-existent sheet '{ref}'")
                        # Flag hardcoded numbers mixed with cell refs
                        has_cell_ref = bool(re.search(r'[A-Z]+\d+', formula_text[1:]))
                        hardcoded_nums = re.findall(r'(?<![A-Z])(\d{2,}(?:\.\d+)?)(?!\d*[A-Z])', formula_text[1:])
                        if has_cell_ref and hardcoded_nums:
                            issues.append(f"hardcoded number(s) {', '.join(hardcoded_nums[:3])} may belong in cells")
                        # Basic circular reference hint (self-reference)
                        if cell.coordinate in formula_text:
                            issues.append("possible circular reference (self-referencing)")

                        formulas_found.append({
                            "sheet": sheet_name,
                            "cell": cell.coordinate,
                            "formula": formula_text,
                            "value": computed,
                            "issues": issues,
                        })
                if len(formulas_found) >= 100:
                    break
        wb_formulas.close()
        wb_values.close()
    except Exception:
        pass  # Formula extraction is best-effort; don't break main parsing

    # Build text summary
    text_parts = []
    for i, table in enumerate(tables):
        sheet_name = sheets[i]["sheet"] if i < len(sheets) else f"Sheet {i+1}"
        text_parts.append(f"## {sheet_name}")
        for row in table[:100]:  # Max 100 rows per sheet
            text_parts.append(" | ".join(row))

    # Append formula summary
    if formulas_found:
        formula_sheets = set(f["sheet"] for f in formulas_found)
        text_parts.append("")
        text_parts.append("## Formulas")
        text_parts.append(f"Formulas found: {len(formulas_found)} formulas across {len(formula_sheets)} sheet(s)")
        text_parts.append("")
        for entry in formulas_found[:20]:
            line = f"  [{entry['sheet']}!{entry['cell']}] {entry['formula']} => {entry['value']}"
            if entry.get("issues"):
                line += f"  ** {'; '.join(entry['issues'])}"
            text_parts.append(line)
        if len(formulas_found) > 20:
            text_parts.append(f"  ... and {len(formulas_found) - 20} more formulas")

    result = {
        "text": "\n".join(text_parts),
        "tables": tables,
        "metadata": {"sheets": [s["sheet"] for s in sheets]},
    }
    if formulas_found:
        result["formulas"] = formulas_found
    return result


def _parse_pptx(filepath):
    try:
        from pptx import Presentation
    except ImportError:
        return {"text": "", "error": "python-pptx not installed"}

    prs = Presentation(filepath)
    slides = []

    for i, slide in enumerate(prs.slides[:50]):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    texts.append(" | ".join(cells))
        if texts:
            slides.append(f"--- Slide {i+1} ---\n" + "\n".join(texts))

    return {
        "text": "\n\n".join(slides),
        "metadata": {"slides": len(prs.slides)},
    }


def _parse_csv(filepath):
    with open(filepath, "r", errors="ignore") as f:
        reader = csv.reader(f)
        rows = []
        for i, row in enumerate(reader):
            if i >= 200:
                break
            rows.append(row)

    text_parts = []
    for row in rows:
        text_parts.append(" | ".join(row))

    return {"text": "\n".join(text_parts), "tables": [rows]}


def _score_section(text: str) -> int:
    """Score a text section by analytical value. Higher = keep first during truncation."""
    t = text.lower()
    score = 50  # baseline

    # High-value: financials, metrics, claims, projections
    high_signals = ["revenue", "profit", "margin", "cost", "budget", "forecast",
                    "projection", "growth", "target", "goal", "kpi", "metric",
                    "valuation", "arpu", "cac", "ltv", "churn", "mrr", "arr",
                    "burn rate", "runway", "roi", "irr", "npv", "ebitda",
                    "total addressable", "market size", "cagr"]
    score += sum(15 for w in high_signals if w in t)

    # Executive summary / key sections
    summary_signals = ["executive summary", "key findings", "recommendation",
                       "conclusion", "overview", "highlights", "key takeaway",
                       "summary", "strategic", "action item", "next step",
                       "risk", "assumption", "critical"]
    score += sum(20 for w in summary_signals if w in t)

    # Numbers = data-rich sections
    import re
    numbers = re.findall(r'\$[\d,.]+[KMBkmb]?|\d+(?:\.\d+)?%|\d{1,3}(?:,\d{3})+', t)
    score += min(len(numbers) * 5, 60)  # cap at 60 bonus for number density

    # Tables (pipe-separated) = structured data
    pipe_rows = sum(1 for line in text.split('\n') if line.count('|') >= 2)
    score += min(pipe_rows * 3, 30)

    # Low-value: boilerplate, legal, appendices
    low_signals = ["disclaimer", "terms and conditions", "all rights reserved",
                   "confidential", "table of contents", "appendix", "lorem ipsum",
                   "copyright", "proprietary", "acknowledgment"]
    score -= sum(20 for w in low_signals if w in t)

    return max(score, 1)


def _smart_truncate(file_texts: list, max_chars: int = 100000) -> str:
    """Truncate intelligently: keep high-value sections, trim boilerplate.

    file_texts: list of (filename, full_text) tuples
    Returns combined text within max_chars budget.
    """
    total = sum(len(t) for _, t in file_texts)
    if total <= max_chars:
        return "\n\n".join(f"=== {name} ===\n{text}" for name, text in file_texts)

    # Split each file into sections, score them, allocate budget
    all_sections = []  # (filename, section_text, score, original_order)
    order = 0
    for filename, text in file_texts:
        # Split on headers, page breaks, or double newlines
        import re
        chunks = re.split(r'(?:\n(?:---|===|#{1,3} )|\n\n--- Page Break ---\n\n|\n\n\n)', text)
        if not chunks or (len(chunks) == 1 and len(text) > 5000):
            # Large unsplit text — chunk by paragraphs (every ~2000 chars)
            chunks = []
            for i in range(0, len(text), 2000):
                end = text.find('\n', i + 1800)
                if end == -1 or end > i + 2500:
                    end = i + 2000
                chunks.append(text[i:end])

        for chunk in chunks:
            chunk = chunk.strip()
            if not chunk:
                continue
            score = _score_section(chunk)
            all_sections.append((filename, chunk, score, order))
            order += 1

    # Sort by score descending, but keep file headers (first section of each file) boosted
    first_sections = set()
    for fn, _, _, o in all_sections:
        if fn not in first_sections:
            first_sections.add(fn)
            # Boost first section of each file (context/overview)
            for i, (f, c, s, oo) in enumerate(all_sections):
                if f == fn and oo == o:
                    all_sections[i] = (f, c, s + 100, oo)
                    break

    # Sort by score desc, take until budget full
    ranked = sorted(all_sections, key=lambda x: -x[2])
    budget = max_chars - 500  # leave room for truncation note
    kept = []
    used = 0
    for filename, chunk, score, orig_order in ranked:
        chunk_size = len(chunk) + len(filename) + 10  # overhead for header
        if used + chunk_size > budget:
            # Try to fit a trimmed version of this section
            remaining = budget - used - len(filename) - 50
            if remaining > 200:
                kept.append((filename, chunk[:remaining] + "\n[Section trimmed]", orig_order))
                used += remaining + len(filename) + 50
            break
        kept.append((filename, chunk, orig_order))
        used += chunk_size

    # Re-sort by original order so document flow is preserved
    kept.sort(key=lambda x: x[2])

    # Group by filename for readable output
    output = []
    current_file = None
    for filename, chunk, _ in kept:
        if filename != current_file:
            output.append(f"\n=== {filename} ===")
            current_file = filename
        output.append(chunk)

    result = "\n\n".join(output)
    dropped = len(all_sections) - len(kept)
    if dropped > 0:
        result += f"\n\n[{dropped} lower-priority sections omitted to fit context window. Key financials, metrics, and executive summaries were preserved.]"
    return result


def _extract_claims(files: list) -> str:
    """Extract numeric claims and key assertions from each file for cross-referencing."""
    import re
    claims_by_file = {}

    for f in files:
        text = f.get("text", "")
        if not text or len(text) < 20:
            continue

        claims = []

        # Dollar amounts with context
        for m in re.finditer(r'(.{0,60}\$[\d,.]+[KMBkmb]?\s*(?:million|billion|thousand|k|m|b)?[^.\n]{0,60})', text, re.IGNORECASE):
            claim = m.group(1).strip()
            if len(claim) > 20:
                claims.append(claim)

        # Percentages with context
        for m in re.finditer(r'(.{0,60}\d+(?:\.\d+)?%[^.\n]{0,60})', text):
            claim = m.group(1).strip()
            if len(claim) > 15:
                claims.append(claim)

        # Growth/revenue/metric claims
        for m in re.finditer(r'(.{0,40}(?:revenue|profit|growth|margin|users|customers|employees|headcount|valuation|funding|raised|arr|mrr|burn)[^.\n]{0,80})', text, re.IGNORECASE):
            claim = m.group(1).strip()
            if len(claim) > 20 and claim not in claims:
                claims.append(claim)

        # Date/timeline claims
        for m in re.finditer(r'(.{0,40}(?:by Q[1-4]|by 20\d\d|launch|deadline|milestone|target date|go.live)[^.\n]{0,60})', text, re.IGNORECASE):
            claim = m.group(1).strip()
            if len(claim) > 15 and claim not in claims:
                claims.append(claim)

        if claims:
            # Deduplicate and limit
            seen = set()
            unique = []
            for c in claims:
                # Clean leading punctuation/whitespace
                clean = c.lstrip(' .,;:-')
                if len(clean) < 15:
                    continue
                short = clean[:50].lower()
                if short not in seen:
                    seen.add(short)
                    unique.append(clean)
            claims_by_file[f["filename"]] = unique[:25]  # max 25 claims per file

    if not claims_by_file:
        return ""

    lines = ["\n## Cross-Reference Map — Key Claims by File"]
    lines.append("Verify these claims are consistent across documents. Flag any contradictions.\n")
    for filename, claims in claims_by_file.items():
        lines.append(f"**{filename}:**")
        for c in claims:
            # Clean up whitespace
            clean = ' '.join(c.split())
            lines.append(f"  - {clean}")
        lines.append("")

    return "\n".join(lines)


def _analyze_code_relationships(files: list) -> str:
    """Analyze relationships between code files."""
    code_exts = {'.py', '.js', '.ts', '.jsx', '.tsx', '.java', '.go', '.rs', '.rb', '.php', '.c', '.cpp', '.h', '.cs', '.swift', '.kt'}

    code_files = [f for f in files if any(f['filename'].endswith(ext) for ext in code_exts)]
    if not code_files:
        return ""

    imports = {}  # filename -> [imported modules]
    definitions = {}  # filename -> [function/class names]
    project_type = "Unknown"

    # Detect project type
    filenames = {f['filename'] for f in files}
    if any('package.json' in fn for fn in filenames): project_type = "Node.js"
    if any('requirements.txt' in fn for fn in filenames): project_type = "Python"
    if any('Cargo.toml' in fn for fn in filenames): project_type = "Rust"
    if any('go.mod' in fn for fn in filenames): project_type = "Go"
    if any('pom.xml' in fn for fn in filenames): project_type = "Java (Maven)"
    if any('Gemfile' in fn for fn in filenames): project_type = "Ruby"
    if any('Dockerfile' in fn for fn in filenames): project_type += " + Docker"
    if any('docker-compose' in fn for fn in filenames): project_type += " + Docker Compose"

    # Detect framework
    for f in files:
        text = f.get('text', '')[:2000]
        if 'from flask' in text.lower() or 'import flask' in text.lower(): project_type += " + Flask"
        if 'from django' in text.lower(): project_type += " + Django"
        if 'import express' in text.lower() or 'require.*express' in text.lower(): project_type += " + Express"
        if 'import react' in text.lower(): project_type += " + React"

    import re

    for f in code_files:
        fname = f['filename']
        text = f.get('text', '')

        # Extract imports
        file_imports = []
        # Python: from X import Y, import X
        for m in re.findall(r'^(?:from\s+(\S+)\s+import|import\s+(\S+))', text, re.MULTILINE):
            mod = m[0] or m[1]
            if mod: file_imports.append(mod.split('.')[0])
        # JS/TS: import X from 'Y', require('Y')
        for m in re.findall(r"(?:import\s+.*?from\s+['\"](.+?)['\"]|require\s*\(\s*['\"](.+?)['\"])", text):
            mod = m[0] or m[1]
            if mod: file_imports.append(mod.split('/')[0].lstrip('.'))
        # Go: import "X"
        for m in re.findall(r'import\s+["\'](.+?)["\']', text):
            file_imports.append(m.split('/')[-1])

        imports[fname] = list(set(file_imports))

        # Extract definitions
        file_defs = []
        # Python: def X, class X
        for m in re.findall(r'^(?:def|class)\s+(\w+)', text, re.MULTILINE):
            file_defs.append(m)
        # JS/TS: function X, const X = (...) =>, class X, export function X
        for m in re.findall(r'^(?:export\s+)?(?:function|class|const|let|var)\s+(\w+)', text, re.MULTILINE):
            file_defs.append(m)
        # Go: func X
        for m in re.findall(r'^func\s+(?:\([^)]*\)\s+)?(\w+)', text, re.MULTILINE):
            file_defs.append(m)

        definitions[fname] = file_defs[:50]  # Cap at 50 per file

    # Build cross-file reference map
    all_defs = {}  # name -> filename
    for fname, defs in definitions.items():
        for d in defs:
            if len(d) > 2:  # Skip very short names
                all_defs[d] = fname

    cross_refs = {}  # filename -> [(called_function, defined_in)]
    for f in code_files:
        fname = f['filename']
        text = f.get('text', '')
        refs = []
        for func_name, defined_in in all_defs.items():
            if defined_in != fname and func_name in text:
                refs.append((func_name, defined_in))
        if refs:
            cross_refs[fname] = refs[:20]  # Cap

    # Build output
    lines = [f"## Code Architecture ({project_type})\n"]
    lines.append(f"**{len(code_files)} code files** across {len(set(f['filename'].split('/')[0] for f in code_files if '/' in f['filename']))} directories\n")

    # Dependency graph
    if imports:
        lines.append("### Dependencies (who imports whom)")
        for fname, imps in sorted(imports.items()):
            if imps:
                # Check if any import matches another uploaded file
                internal = [i for i in imps if any(i in other['filename'] for other in code_files)]
                external = [i for i in imps if i not in internal]
                if internal:
                    lines.append(f"- **{fname}** \u2192 internal: {', '.join(internal)}")
                if external and len(external) <= 10:
                    lines.append(f"  external: {', '.join(external[:10])}")
        lines.append("")

    # Key definitions
    if definitions:
        lines.append("### Key Definitions (functions/classes per file)")
        for fname, defs in sorted(definitions.items()):
            if defs:
                lines.append(f"- **{fname}**: {', '.join(defs[:10])}" + (f" (+{len(defs)-10} more)" if len(defs) > 10 else ""))
        lines.append("")

    # Cross-file calls
    if cross_refs:
        lines.append("### Cross-File References (who calls what from where)")
        for fname, refs in sorted(cross_refs.items()):
            ref_strs = [f"`{fn}` (from {src})" for fn, src in refs[:5]]
            lines.append(f"- **{fname}** uses: {', '.join(ref_strs)}")
        lines.append("")

    return "\n".join(lines)


def parse_folder(filepaths: list) -> dict:
    """Parse multiple files and return combined context with smart prioritization."""
    results = []
    file_texts = []

    for fp in filepaths:
        result = parse_file(fp)
        results.append(result)
        if result["text"]:
            file_texts.append((result["filename"], result["text"]))

    combined_text = _smart_truncate(file_texts, max_chars=100000)
    claims_map = _extract_claims(results)
    code_architecture = _analyze_code_relationships(results)

    # Prepend code architecture so the AI sees structure before raw code
    if code_architecture:
        combined_text = code_architecture + "\n\n" + combined_text

    return {
        "files": results,
        "file_count": len(results),
        "total_chars": len(combined_text),
        "combined_text": combined_text,
        "claims_map": claims_map,
        "code_architecture": code_architecture,
        "errors": [r for r in results if r.get("error")],
    }
